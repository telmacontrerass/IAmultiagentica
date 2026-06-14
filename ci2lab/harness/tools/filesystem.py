"""Herramientas de lectura y listado de archivos."""

from __future__ import annotations

import csv
import io
import logging
import re
from pathlib import Path

from ci2lab.harness.tools.paths import PathViolationError, format_size, resolve_path
from ci2lab.harness.tools.secret_files import (
    grep_skip_notice,
    is_sensitive_path,
    secret_file_block_message,
)


def _resolve_or_error(raw: str, cwd: str) -> tuple[Path | None, str | None]:
    try:
        return resolve_path(raw, cwd), None
    except PathViolationError as exc:
        return None, f"Error: {exc}"


def _resolve_for_access(
    raw: str,
    cwd: str,
    *,
    security_engine: str = "ci2lab",
) -> tuple[Path | None, str | None]:
    from ci2lab.security.engine import enforce_ci2lab_hard_policy

    if enforce_ci2lab_hard_policy(security_engine):
        return _resolve_or_error(raw, cwd)
    try:
        candidate = Path(raw).expanduser()
        if not candidate.is_absolute():
            candidate = Path(cwd) / candidate
        return candidate.resolve(), None
    except OSError as exc:
        return None, f"Error: ruta invalida: {exc}"


def _check_sensitive(
    resolved: Path,
    cwd: str,
    *,
    security_engine: str = "ci2lab",
) -> bool:
    from ci2lab.security.engine import enforce_ci2lab_hard_policy

    if not enforce_ci2lab_hard_policy(security_engine):
        return False
    return is_sensitive_path(resolved, workspace=cwd)

MAX_READ_LINES = 2000
MAX_PDF_PAGES = 100
MAX_SPREADSHEET_ROWS = 200
MAX_SPREADSHEET_COLS = 30
MAX_PRESENTATION_SLIDES = 200
MAX_DOCUMENT_CHARS = 120_000

TEXT_DOCUMENT_SUFFIXES = frozenset({
    ".csv",
    ".json",
    ".log",
    ".md",
    ".rst",
    ".text",
    ".tsv",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
})
OFFICE_DOCUMENT_SUFFIXES = frozenset({".docx", ".pptx", ".xlsx"})
SUPPORTED_DOCUMENT_SUFFIXES = (
    TEXT_DOCUMENT_SUFFIXES | OFFICE_DOCUMENT_SUFFIXES | frozenset({".pdf"})
)


def read_file(
    cwd: str,
    path: str,
    offset: int = 1,
    limit: int | None = None,
    *,
    security_engine: str = "ci2lab",
) -> str:
    resolved, err = _resolve_for_access(path, cwd, security_engine=security_engine)
    if err:
        return err
    assert resolved is not None
    if not resolved.is_file():
        return f"Error: no existe el archivo {resolved}"
    if _check_sensitive(resolved, cwd, security_engine=security_engine):
        return secret_file_block_message()
    text = extract_document_text(resolved, include_metadata=False)
    if text.startswith("Error:"):
        return text
    return _numbered_lines(text, offset=offset, limit=limit)


def read_document(cwd: str, path: str) -> str:
    """Read a document-like file and return structured extracted text."""
    resolved, err = _resolve_or_error(path, cwd)
    if err:
        return err
    assert resolved is not None
    if not resolved.is_file():
        return f"Error: no existe el archivo {resolved}"
    if is_sensitive_path(resolved):
        return secret_file_block_message()
    return extract_document_text(resolved, include_metadata=True)


def _numbered_lines(text: str, offset: int = 1, limit: int | None = None) -> str:
    lines = text.splitlines()
    start = max(1, offset if offset is not None else 1)
    end = start + (limit or MAX_READ_LINES) - 1
    slice_lines = lines[start - 1 : end]
    numbered = [f"{i + start:6d}|{line}" for i, line in enumerate(slice_lines)]
    if len(lines) > end:
        numbered.append(f"... ({len(lines) - end} líneas más; usa offset/limit)")
    return "\n".join(numbered) if numbered else "(archivo vacío)"


def extract_document_text(path: Path, *, include_metadata: bool = False) -> str:
    """Extract text from supported teaching/document formats."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        text = extract_pdf_text(path)
        sections = (
            _pdf_section_count(path)
            if include_metadata and not text.startswith("Error:")
            else None
        )
    elif suffix == ".docx":
        text = extract_docx_text(path)
        sections = "desconocido"
    elif suffix == ".pptx":
        text, sections = extract_pptx_text(path)
    elif suffix == ".xlsx":
        text, sections = extract_xlsx_text(path)
    elif suffix in {".csv", ".tsv"}:
        text = extract_csv_text(path)
        sections = "1 tabla"
    elif suffix in TEXT_DOCUMENT_SUFFIXES or not suffix:
        text = path.read_text(encoding="utf-8", errors="replace")
        sections = "texto plano"
    elif not include_metadata:
        text = path.read_text(encoding="utf-8", errors="replace")
        sections = "texto plano"
    else:
        return (
            f"Error: formato no soportado para lectura documental: "
            f"{suffix or '(sin extension)'}"
        )

    if text.startswith("Error:") or not include_metadata:
        return text

    text = _truncate_document_text(text)
    return (
        f"Documento: {path.name}\n"
        f"Tipo: {suffix.lstrip('.') or 'texto'}\n"
        f"Paginas/secciones: {sections or 'desconocido'}\n"
        "Texto extraido:\n\n"
        f"{text}"
    ).strip()


def _truncate_document_text(text: str) -> str:
    if len(text) <= MAX_DOCUMENT_CHARS:
        return text
    return (
        text[:MAX_DOCUMENT_CHARS].rstrip()
        + f"\n\n... (texto truncado; limite {MAX_DOCUMENT_CHARS} caracteres)"
    )


def _pdf_section_count(path: Path) -> str | None:
    try:
        from pypdf import PdfReader

        return f"{len(PdfReader(str(path)).pages)} paginas"
    except Exception:  # noqa: BLE001
        return None


def extract_pdf_text(path: Path) -> str:
    """Extract text from a PDF file, returning an Error: string on failure."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return (
            "Error: no se puede leer PDF porque falta la dependencia `pypdf`. "
            "Instala el proyecto de nuevo para activar soporte PDF."
        )

    logging.getLogger("pypdf").setLevel(logging.ERROR)

    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001
        return f"Error: no se pudo abrir el PDF {path}: {exc}"

    total_pages = len(reader.pages)
    page_count = min(total_pages, MAX_PDF_PAGES)
    chunks: list[str] = []
    has_extractable_text = False
    for index in range(page_count):
        page = reader.pages[index]
        page_number = index + 1
        chunks.append(f"[PDF page {page_number}/{total_pages}]")
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            page_text = f"Error extrayendo texto de esta pagina: {exc}"
        page_text = page_text.strip()
        if page_text:
            has_extractable_text = True
        chunks.append(page_text or "(sin texto extraible en esta pagina)")

    if total_pages > page_count:
        chunks.append(
            f"... ({total_pages - page_count} paginas mas; limite PDF {MAX_PDF_PAGES})"
        )

    if not has_extractable_text:
        return (
            "Error: el PDF no contiene texto extraible. Puede ser un PDF escaneado; "
            "hace falta OCR para leer imagenes."
        )
    return "\n".join(chunks).strip()


def extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return (
            "Error: no se puede leer DOCX porque falta la dependencia `python-docx`. "
            "Instala el proyecto de nuevo para activar soporte Word."
        )

    try:
        document = Document(str(path))
    except Exception as exc:  # noqa: BLE001
        return f"Error: no se pudo abrir el DOCX {path}: {exc}"

    chunks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
        if style_name.lower().startswith("heading"):
            chunks.append(f"[{style_name}] {text}")
        else:
            chunks.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            chunks.append(f"[Table {table_index}]")
            chunks.extend(rows)

    return "\n".join(chunks).strip() or "(documento DOCX sin texto extraible)"


def extract_pptx_text(path: Path) -> tuple[str, str]:
    try:
        from pptx import Presentation
    except ImportError:
        return (
            "Error: no se puede leer PPTX porque falta la dependencia `python-pptx`. "
            "Instala el proyecto de nuevo para activar soporte PowerPoint."
        ), "desconocido"

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        return f"Error: no se pudo abrir el PPTX {path}: {exc}", "desconocido"

    total_slides = len(presentation.slides)
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > MAX_PRESENTATION_SLIDES:
            break
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text).strip()
                if text:
                    slide_text.append(text)
        chunks.append(f"[Slide {slide_index}]")
        chunks.append("\n".join(slide_text) if slide_text else "(sin texto extraible)")

    if total_slides > MAX_PRESENTATION_SLIDES:
        chunks.append(
            f"... ({total_slides - MAX_PRESENTATION_SLIDES} diapositivas mas; "
            f"limite {MAX_PRESENTATION_SLIDES})"
        )

    return "\n".join(chunks).strip(), f"{total_slides} diapositivas"


def extract_xlsx_text(path: Path) -> tuple[str, str]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return (
            "Error: no se puede leer XLSX porque falta la dependencia `openpyxl`. "
            "Instala el proyecto de nuevo para activar soporte Excel."
        ), "desconocido"

    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        return f"Error: no se pudo abrir el XLSX {path}: {exc}", "desconocido"

    chunks: list[str] = []
    for sheet in workbook.worksheets:
        chunks.append(f"[Sheet: {sheet.title}]")
        row_count = 0
        for row in sheet.iter_rows(
            max_row=MAX_SPREADSHEET_ROWS,
            max_col=MAX_SPREADSHEET_COLS,
            values_only=True,
        ):
            row_count += 1
            values = _trim_empty_tail(["" if value is None else str(value) for value in row])
            if not values:
                continue
            chunks.append(" | ".join(values).rstrip())
        if row_count >= MAX_SPREADSHEET_ROWS:
            chunks.append(
                f"... (hoja truncada; limite {MAX_SPREADSHEET_ROWS} filas x "
                f"{MAX_SPREADSHEET_COLS} columnas)"
            )

    close = getattr(workbook, "close", None)
    if callable(close):
        close()
    return "\n".join(chunks).strip(), f"{len(workbook.worksheets)} hojas"


def extract_csv_text(path: Path) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    delimiter = "\t" if path.suffix.lower() == ".tsv" else None
    sample = text[:4096]
    if delimiter is None:
        try:
            delimiter = csv.Sniffer().sniff(sample).delimiter
        except csv.Error:
            delimiter = ","
    rows: list[str] = []
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    for index, row in enumerate(reader, start=1):
        if index > MAX_SPREADSHEET_ROWS:
            rows.append(f"... (tabla truncada; limite {MAX_SPREADSHEET_ROWS} filas)")
            break
        values = _trim_empty_tail(row)
        if values:
            rows.append(" | ".join(values))
    return "\n".join(rows).strip() or "(archivo CSV vacio)"


def _trim_empty_tail(values: list[str]) -> list[str]:
    trimmed = [value.strip() for value in values]
    while trimmed and not trimmed[-1]:
        trimmed.pop()
    return trimmed


def ls(cwd: str, path: str = ".") -> str:
    resolved, err = _resolve_or_error(path or ".", cwd)
    if err:
        return err
    assert resolved is not None
    if not resolved.is_dir():
        return f"Error: no es un directorio {resolved}"
    dirs: list[str] = []
    files: list[str] = []
    for entry in sorted(resolved.iterdir(), key=lambda p: p.name.lower()):
        if entry.name.startswith("."):
            continue
        if entry.is_dir():
            dirs.append(f"  {entry.name}/")
        elif entry.is_file():
            files.append(f"  {entry.name}  ({format_size(entry.stat().st_size)})")
    lines = [f"{resolved}/"]
    lines.extend(dirs)
    lines.extend(files)
    return "\n".join(lines) if len(lines) > 1 else f"{resolved}/ (vacío)"


def glob_search(cwd: str, pattern: str, path: str = ".") -> str:
    base, err = _resolve_or_error(path or ".", cwd)
    if err:
        return err
    assert base is not None
    if not base.is_dir():
        return f"Error: base no es directorio {base}"
    matches = sorted(base.glob(pattern), key=lambda p: p.stat().st_mtime, reverse=True)
    if not matches:
        return f"Sin coincidencias para `{pattern}` en {base}"
    lines = [str(m.relative_to(Path(cwd).resolve())) for m in matches[:100]]
    if len(matches) > 100:
        lines.append(f"... y {len(matches) - 100} más")
    return "\n".join(lines)


def grep_search(
    cwd: str,
    pattern: str,
    path: str = ".",
    glob_pattern: str | None = None,
    ignore_case: bool = False,
    max_results: int = 50,
) -> str:
    base, err = _resolve_or_error(path or ".", cwd)
    if err:
        return err
    assert base is not None
    flags = re.IGNORECASE if ignore_case else 0
    try:
        regex = re.compile(pattern, flags)
    except re.error as exc:
        return f"Error: expresión regular inválida: {exc}"

    if base.is_file():
        if is_sensitive_path(base, workspace=cwd):
            return secret_file_block_message()
        return _grep_single_file(base, root=Path(cwd).resolve(), regex=regex, max_results=max_results)

    results, skipped = _grep_scan_tree(
        base,
        root=Path(cwd).resolve(),
        regex=regex,
        glob_pattern=glob_pattern,
        max_results=max_results,
    )
    if results:
        body = "\n".join(results)
        notice = grep_skip_notice(skipped)
        return f"{body}\n{notice}" if notice else body
    if skipped:
        notice = grep_skip_notice(skipped)
        return f"Sin coincidencias para `{pattern}`\n{notice}"
    return f"Sin coincidencias para `{pattern}`"


def _grep_single_file(
    file_path: Path,
    *,
    root: Path,
    regex: re.Pattern[str],
    max_results: int,
) -> str:
    results: list[str] = []
    try:
        rel = file_path.relative_to(root)
    except ValueError:
        rel = file_path
    try:
        content = extract_document_text(file_path, include_metadata=False)
    except OSError:
        return f"Sin coincidencias para `{regex.pattern}`"
    if content.startswith("Error:"):
        return f"Sin coincidencias para `{regex.pattern}`"
    for i, line in enumerate(content.splitlines(), start=1):
        if regex.search(line):
            results.append(f"{rel}:{i}:{line}")
            if len(results) >= max_results:
                break
    return "\n".join(results) if results else f"Sin coincidencias para `{regex.pattern}`"


def _grep_scan_tree(
    base: Path,
    *,
    root: Path,
    regex: re.Pattern[str],
    glob_pattern: str | None,
    max_results: int,
) -> tuple[list[str], int]:
    results: list[str] = []
    skipped = 0
    for file_path in base.rglob("*"):
        if not file_path.is_file():
            continue
        if glob_pattern and not file_path.match(glob_pattern):
            continue
        if is_sensitive_path(file_path, workspace=root):
            skipped += 1
            continue
        try:
            rel = file_path.relative_to(root)
        except ValueError:
            continue
        try:
            content = extract_document_text(file_path, include_metadata=False)
        except OSError:
            continue
        if content.startswith("Error:"):
            continue
        for i, line in enumerate(content.splitlines(), start=1):
            if regex.search(line):
                results.append(f"{rel}:{i}:{line}")
                if len(results) >= max_results:
                    return results, skipped
    return results, skipped


def write_file(cwd: str, path: str, content: str) -> str:
    resolved, err = _resolve_or_error(path, cwd)
    if err:
        return err
    assert resolved is not None
    if is_sensitive_path(resolved, workspace=cwd):
        return secret_file_block_message()
    resolved.parent.mkdir(parents=True, exist_ok=True)
    resolved.write_text(content, encoding="utf-8")
    return f"Escrito {resolved} ({len(content)} caracteres)"


def edit_file(
    cwd: str,
    path: str,
    old_string: str,
    new_string: str,
    replace_all: bool = False,
) -> str:
    from ci2lab.harness.tools.write_preview import compute_edit_result

    resolved, err = _resolve_or_error(path, cwd)
    if err:
        return err
    assert resolved is not None
    if is_sensitive_path(resolved, workspace=cwd):
        return secret_file_block_message()
    if resolved.is_file():
        original_count = resolved.read_text(encoding="utf-8", errors="replace").count(
            old_string
        )
    else:
        original_count = 0

    new_text, error = compute_edit_result(
        cwd, path, old_string, new_string, replace_all
    )
    if error:
        return error
    resolved.write_text(new_text or "", encoding="utf-8")
    replaced = original_count if replace_all else 1
    return f"Editado {resolved}: {replaced} reemplazo(s)"


def permission_summary(tool_name: str, args: dict) -> str:
    """Resumen corto para el diálogo de confirmación."""
    if tool_name == "bash":
        cmd = args.get("command", "")
        return cmd[:120] + ("..." if len(cmd) > 120 else "")
    if tool_name in ("write_file", "edit_file", "notebook_edit"):
        return str(args.get("path", ""))
    if tool_name == "apply_patch":
        patch = str(args.get("patch", ""))
        return patch[:120] + ("..." if len(patch) > 120 else "")
    if tool_name == "fill_docx_template":
        return f"{args.get('template', '')} → {args.get('output', '')}"
    if tool_name == "web_fetch":
        return str(args.get("url", ""))[:120]
    if tool_name == "ask_user":
        return str(args.get("question", ""))[:120]
    return str(args)[:80]
