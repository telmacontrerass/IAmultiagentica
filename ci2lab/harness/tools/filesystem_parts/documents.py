"""Document readers and text extraction helpers."""

from __future__ import annotations

import csv
import io
import logging
from pathlib import Path
from types import ModuleType

#: Maximum number of lines returned by a single ``read_file`` call.
MAX_READ_LINES: int = 2000
#: Maximum number of PDF pages whose text is extracted.
MAX_PDF_PAGES: int = 100
#: Maximum number of spreadsheet rows read per sheet.
MAX_SPREADSHEET_ROWS: int = 200
#: Maximum number of spreadsheet columns read per row.
MAX_SPREADSHEET_COLS: int = 30
#: Maximum number of presentation slides whose text is extracted.
MAX_PRESENTATION_SLIDES: int = 200
#: Maximum number of characters returned for an extracted document.
MAX_DOCUMENT_CHARS: int = 120_000

#: Plain-text document suffixes read directly as UTF-8.
TEXT_DOCUMENT_SUFFIXES: frozenset[str] = frozenset(
    {
        ".csv",
        ".json",
        ".log",
        ".md",
        ".rst",
        ".text",
        ".tsv",
        ".txt",
        ".xml",
        ".yaml",
        ".yml",
    }
)
#: Office (OOXML) document suffixes requiring a dedicated extractor.
OFFICE_DOCUMENT_SUFFIXES: frozenset[str] = frozenset({".docx", ".pptx", ".xlsx"})
#: All document suffixes supported by :func:`extract_document_text`.
SUPPORTED_DOCUMENT_SUFFIXES: frozenset[str] = (
    TEXT_DOCUMENT_SUFFIXES | OFFICE_DOCUMENT_SUFFIXES | frozenset({".pdf"})
)


def numbered_lines(text: str, offset: int = 1, limit: int | None = None) -> str:
    """Render ``text`` as a numbered, line-bounded slice.

    Args:
        text: The full text to slice.
        offset: 1-based line number to start from (clamped to at least 1).
        limit: Maximum number of lines to include, or ``None`` for
            :data:`MAX_READ_LINES`.

    Returns:
        The numbered slice (``"<num>|<line>"`` per row) with a trailing
        truncation hint when more lines follow, or ``"(empty file)"``.
    """
    lines = text.splitlines()
    start = max(1, offset if offset is not None else 1)
    end = start + (limit or MAX_READ_LINES) - 1
    slice_lines = lines[start - 1 : end]
    numbered = [f"{i + start:6d}|{line}" for i, line in enumerate(slice_lines)]
    if len(lines) > end:
        numbered.append(f"... ({len(lines) - end} more lines; use offset/limit)")
    return "\n".join(numbered) if numbered else "(empty file)"


def extract_document_text(path: Path, *, include_metadata: bool = False) -> str:
    """Extract text from supported teaching/document formats.

    Dispatches on the file suffix to the appropriate extractor (PDF, DOCX,
    PPTX, XLSX, CSV/TSV, or plain text). Files without a recognised suffix are
    read as plain text unless ``include_metadata`` requires a known format.

    Args:
        path: The document to read.
        include_metadata: When ``True``, prepend a metadata header (name, type,
            page/section count) and truncate to :data:`MAX_DOCUMENT_CHARS`.

    Returns:
        The extracted text (optionally with a metadata header), or an
        ``"Error: ..."`` message for unsupported formats or extraction errors.
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        facade = _filesystem_facade()
        text = facade.extract_pdf_text(path)
        sections = (
            facade._pdf_section_count(path)
            if include_metadata and not text.startswith("Error:")
            else None
        )
    elif suffix == ".docx":
        text = extract_docx_text(path)
        sections = "unknown"
    elif suffix == ".pptx":
        text, sections = extract_pptx_text(path)
    elif suffix == ".xlsx":
        text, sections = extract_xlsx_text(path)
    elif suffix in {".csv", ".tsv"}:
        text = extract_csv_text(path)
        sections = "1 table"
    elif suffix in TEXT_DOCUMENT_SUFFIXES or not suffix or not include_metadata:
        text = path.read_text(encoding="utf-8", errors="replace")
        sections = "plain text"
    else:
        return f"Error: unsupported format for document reading: {suffix or '(no extension)'}"

    if text.startswith("Error:") or not include_metadata:
        return text

    text = truncate_document_text(text)
    return (
        f"Document: {path.name}\n"
        f"Type: {suffix.lstrip('.') or 'text'}\n"
        f"Pages/sections: {sections or 'unknown'}\n"
        "Extracted text:\n\n"
        f"{text}"
    ).strip()


def truncate_document_text(text: str) -> str:
    """Truncate ``text`` to :data:`MAX_DOCUMENT_CHARS`, appending a notice.

    Args:
        text: The extracted document text.

    Returns:
        The original text when within the limit, otherwise a truncated copy
        with a trailing truncation notice.
    """
    if len(text) <= MAX_DOCUMENT_CHARS:
        return text
    return (
        text[:MAX_DOCUMENT_CHARS].rstrip()
        + f"\n\n... (text truncated; limit {MAX_DOCUMENT_CHARS} characters)"
    )


def pdf_section_count(path: Path) -> str | None:
    """Return a ``"<n> pages"`` label for a PDF, or ``None`` on any failure.

    Args:
        path: The PDF file to inspect.

    Returns:
        A human-readable page-count label, or ``None`` if ``pypdf`` is missing
        or the file cannot be read.
    """
    try:
        from pypdf import PdfReader

        return f"{len(PdfReader(str(path)).pages)} pages"
    except Exception:
        return None


def extract_pdf_text(path: Path) -> str:
    """Extract text from a PDF file, returning an Error: string on failure.

    Reads up to :data:`MAX_PDF_PAGES` pages, labelling each and noting any
    remaining pages.

    Args:
        path: The PDF file to read.

    Returns:
        The page-labelled extracted text, or an ``"Error: ..."`` message when
        ``pypdf`` is missing, the file cannot be opened, or no text is
        extractable (e.g. a scanned PDF).
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        return (
            "Error: cannot read PDF because the `pypdf` dependency is missing. "
            "Reinstall the project to enable PDF support."
        )

    logging.getLogger("pypdf").setLevel(logging.ERROR)

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        return f"Error: could not open the PDF {path}: {exc}"

    total_pages = len(reader.pages)
    page_count = min(total_pages, MAX_PDF_PAGES)
    chunks: list[str] = []
    has_extractable_text = False
    for index in range(page_count):
        page = reader.pages[index]
        page_number = index + 1
        chunks.append(f"[PDF page {page_number}/{total_pages}]")
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            page_text = f"Error extracting text from this page: {exc}"
        page_text = page_text.strip()
        if page_text:
            has_extractable_text = True
        chunks.append(page_text or "(no extractable text on this page)")

    if total_pages > page_count:
        chunks.append(f"... ({total_pages - page_count} more pages; PDF limit {MAX_PDF_PAGES})")

    if not has_extractable_text:
        return (
            "Error: the PDF has no extractable text. It may be a scanned PDF; "
            "OCR is needed to read images."
        )
    return "\n".join(chunks).strip()


def pdf_has_extractable_text(
    path: Path | str,
    *,
    min_chars: int = 40,
    max_pages: int = 5,
) -> bool:
    """Return True when the PDF has enough embedded text for ``read_document``.

    Args:
        path: The PDF file to inspect.
        min_chars: Total stripped characters required to consider the PDF
            text-bearing.
        max_pages: Maximum number of leading pages to sample.

    Returns:
        ``True`` once the sampled pages accumulate at least ``min_chars``
        characters; ``False`` if ``pypdf`` is missing or the file is unreadable.
    """
    path = Path(path)

    try:
        from pypdf import PdfReader
    except ImportError:
        return False

    logging.getLogger("pypdf").setLevel(logging.ERROR)

    try:
        reader = PdfReader(str(path))
    except Exception:
        return False

    total_chars = 0
    for index in range(min(len(reader.pages), max_pages)):
        try:
            page_text = (reader.pages[index].extract_text() or "").strip()
        except Exception:
            continue
        total_chars += len(page_text)
        if total_chars >= min_chars:
            return True
    return False


def pdf_needs_vision(path: Path | str) -> bool:
    """Return True for scanned/image-only PDFs that need a vision model.

    Args:
        path: The candidate file to inspect.

    Returns:
        ``True`` when ``path`` is a PDF with no embedded extractable text.
    """
    path = Path(path)
    return path.suffix.lower() == ".pdf" and not pdf_has_extractable_text(path)


def extract_docx_text(path: Path) -> str:
    """Extract paragraph and table text from a DOCX file.

    Headings are prefixed with their style name and tables are rendered as
    pipe-separated rows. Falls back to a Markdown extractor when ``python-docx``
    is unavailable.

    Args:
        path: The DOCX file to read.

    Returns:
        The extracted text, a placeholder when the document has no text, or an
        ``"Error: ..."`` message when the file cannot be opened.
    """
    try:
        from docx import Document
    except ImportError:
        from ci2lab.harness.tools.docx import extract_docx_markdown

        return extract_docx_markdown(path)

    try:
        document = Document(str(path))
    except Exception as exc:
        return f"Error: could not open the DOCX {path}: {exc}"

    chunks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
        if style_name.lower().startswith("heading"):
            chunks.append(f"[{style_name}] {text}")
        else:
            chunks.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            chunks.append(f"[Table {table_index}]")
            chunks.extend(rows)

    return "\n".join(chunks).strip() or "(DOCX document with no extractable text)"


def extract_pptx_text(path: Path) -> tuple[str, str]:
    """Extract per-slide text from a PPTX presentation.

    Reads up to :data:`MAX_PRESENTATION_SLIDES` slides, labelling each and
    noting any remaining slides.

    Args:
        path: The PPTX file to read.

    Returns:
        A ``(text, sections)`` pair where ``text`` is the slide-labelled text
        (or an ``"Error: ..."`` message) and ``sections`` is a slide-count
        label (``"desconocido"`` on failure).
    """
    try:
        from pptx import Presentation
    except ImportError:
        return (
            "Error: cannot read PPTX because the `python-pptx` dependency is missing. "
            "Reinstall the project to enable PowerPoint support."
        ), "desconocido"

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        return f"Error: could not open the PPTX {path}: {exc}", "desconocido"

    total_slides = len(presentation.slides)
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > MAX_PRESENTATION_SLIDES:
            break
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text).strip()
                if text:
                    slide_text.append(text)
        chunks.append(f"[Slide {slide_index}]")
        chunks.append("\n".join(slide_text) if slide_text else "(no extractable text)")

    if total_slides > MAX_PRESENTATION_SLIDES:
        chunks.append(
            f"... ({total_slides - MAX_PRESENTATION_SLIDES} more slides; "
            f"limite {MAX_PRESENTATION_SLIDES})"
        )

    return "\n".join(chunks).strip(), f"{total_slides} slides"


def extract_xlsx_text(path: Path) -> tuple[str, str]:
    """Extract cell text from an XLSX workbook, one section per sheet.

    Reads up to :data:`MAX_SPREADSHEET_ROWS` rows and
    :data:`MAX_SPREADSHEET_COLS` columns per sheet, rendering rows as
    pipe-separated values.

    Args:
        path: The XLSX file to read.

    Returns:
        A ``(text, sections)`` pair where ``text`` is the sheet-labelled text
        (or an ``"Error: ..."`` message) and ``sections`` is a sheet-count
        label (``"desconocido"`` on failure).
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        return (
            "Error: cannot read XLSX because the `openpyxl` dependency is missing. "
            "Reinstall the project to enable Excel support."
        ), "desconocido"

    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        return f"Error: could not open the XLSX {path}: {exc}", "desconocido"

    chunks: list[str] = []
    for sheet in workbook.worksheets:
        chunks.append(f"[Sheet: {sheet.title}]")
        row_count = 0
        for row in sheet.iter_rows(
            max_row=MAX_SPREADSHEET_ROWS,
            max_col=MAX_SPREADSHEET_COLS,
            values_only=True,
        ):
            row_count += 1
            values = trim_empty_tail(["" if value is None else str(value) for value in row])
            if not values:
                continue
            chunks.append(" | ".join(values).rstrip())
        if row_count >= MAX_SPREADSHEET_ROWS:
            chunks.append(
                f"... (sheet truncated; limit {MAX_SPREADSHEET_ROWS} rows x "
                f"{MAX_SPREADSHEET_COLS} columns)"
            )

    close = getattr(workbook, "close", None)
    if callable(close):
        close()
    return "\n".join(chunks).strip(), f"{len(workbook.worksheets)} sheets"


def extract_csv_text(path: Path) -> str:
    """Extract rows from a CSV/TSV file as pipe-separated text.

    The delimiter is the tab for ``.tsv`` files and otherwise sniffed (falling
    back to a comma). Reads up to :data:`MAX_SPREADSHEET_ROWS` rows.

    Args:
        path: The CSV or TSV file to read.

    Returns:
        The pipe-separated rows joined by newlines, or ``"(empty CSV file)"``
        when no data rows are present.
    """
    text = path.read_text(encoding="utf-8", errors="replace")
    delimiter = "\t" if path.suffix.lower() == ".tsv" else None
    sample = text[:4096]
    if delimiter is None:
        try:
            delimiter = csv.Sniffer().sniff(sample).delimiter
        except csv.Error:
            delimiter = ","
    rows: list[str] = []
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    for index, row in enumerate(reader, start=1):
        if index > MAX_SPREADSHEET_ROWS:
            rows.append(f"... (tabla truncada; limite {MAX_SPREADSHEET_ROWS} filas)")
            break
        values = trim_empty_tail(row)
        if values:
            rows.append(" | ".join(values))
    return "\n".join(rows).strip() or "(empty CSV file)"


def trim_empty_tail(values: list[str]) -> list[str]:
    """Strip each value and drop trailing empty entries.

    Args:
        values: The raw cell values for a row.

    Returns:
        The stripped values with empty trailing entries removed.
    """
    trimmed = [value.strip() for value in values]
    while trimmed and not trimmed[-1]:
        trimmed.pop()
    return trimmed


def _filesystem_facade() -> ModuleType:
    """Return the ``filesystem`` facade module (late import to avoid cycles)."""
    from ci2lab.harness.tools import filesystem

    return filesystem
