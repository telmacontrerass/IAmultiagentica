"""PowerPoint writer tool backed by python-pptx."""

from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from ci2lab.harness.tools.paths import PathViolationError, resolve_path
from ci2lab.harness.tools.write_preview import WritePreview

MAX_SLIDES = 40
MAX_TITLE_CHARS = 160
MAX_TEXT_CHARS = 1200
MAX_BULLETS = 12
MAX_BULLET_CHARS = 240
MAX_TABLE_ROWS = 20
MAX_TABLE_COLS = 6
MAX_CELL_CHARS = 160
MAX_CARDS = 6
MAX_FONT_FAMILY_CHARS = 80
MAX_FOOTER_CHARS = 120
MIN_TITLE_FONT_SIZE = 18
MAX_TITLE_FONT_SIZE = 60
MIN_BODY_FONT_SIZE = 8
MAX_BODY_FONT_SIZE = 32
MIN_FITTED_BODY_FONT_SIZE = 12
MIN_COVER_TITLE_FONT_SIZE = 30
MIN_COVER_SUBTITLE_FONT_SIZE = 14

HARD_TEXT_CHARS = 4000
HARD_LIST_ITEM_CHARS = 1200
HARD_CELL_CHARS = 600

COVER_TITLE_LIMIT = 90
COVER_SUBTITLE_LIMIT = 150
TITLE_LAYOUT_LIMIT = 90
BULLET_LAYOUT_LIMIT = 150
COLUMN_ITEM_LIMIT = 115
CARD_LABEL_LIMIT = 34
CARD_VALUE_LIMIT = 26
DECISION_RECOMMENDATION_LIMIT = 160
DECISION_RATIONALE_LIMIT = 220
TABLE_CELL_LAYOUT_LIMIT = 70

SUPPORTED_SLIDE_TYPES = frozenset(
    {
        "cover",
        "section",
        "bullets",
        "two_columns",
        "table",
        "quote",
        "closing",
        "metric_cards",
        "comparison",
        "decision",
    }
)

_HEX_COLOR_RE = re.compile(r"^#?[0-9a-fA-F]{6}$")


@dataclass(frozen=True)
class PresentationPlan:
    """Validated write_pptx request."""

    output_path: Path
    relative_output_path: str
    title: str
    slides: list[dict[str, Any]]
    theme: dict[str, Any]
    expected_titles: list[str]
    overwrite: bool
    layout_warnings: list[dict[str, Any]]


def preview_write_pptx(
    cwd: str,
    output_path: str,
    title: str,
    slides: list[dict[str, Any]],
    theme: dict[str, Any] | None = None,
    overwrite: bool = False,
) -> WritePreview:
    """Validate and preview a PowerPoint write."""
    plan, error = validate_pptx_request(
        cwd,
        output_path=output_path,
        title=title,
        slides=slides,
        theme=theme,
        overwrite=overwrite,
    )
    if error:
        return WritePreview(
            path=output_path or "(no output)",
            is_new_file=True,
            diff="",
            validation_error=error,
        )
    assert plan is not None
    overwrite_note = "existing file; overwrite=true" if plan.output_path.is_file() else "new file"
    lines = [
        f"Output : {plan.relative_output_path} ({overwrite_note})",
        f"Title  : {plan.title}",
        f"Slides : {len(plan.slides)}",
        "Types  : " + ", ".join(str(slide.get("type")) for slide in plan.slides),
    ]
    return WritePreview(
        path=plan.relative_output_path,
        is_new_file=not plan.output_path.is_file(),
        diff="",
        new_content="\n".join(lines),
    )


def validate_pptx_request(
    cwd: str,
    *,
    output_path: str,
    title: str,
    slides: list[dict[str, Any]],
    theme: dict[str, Any] | None = None,
    overwrite: bool = False,
) -> tuple[PresentationPlan | None, str | None]:
    """Validate write_pptx arguments without creating a file."""
    if not isinstance(output_path, str) or not output_path.strip():
        return None, "Error: output_path is required"
    try:
        resolved = resolve_path(output_path, cwd)
    except PathViolationError as exc:
        return None, f"Error: {exc}"

    if resolved.suffix.lower() != ".pptx":
        return None, "Error: write_pptx only accepts .pptx output paths"
    if resolved.exists() and not overwrite:
        return None, "Error: output file already exists; set overwrite=true to replace it"
    if not isinstance(title, str) or not title.strip():
        return None, "Error: title is required"
    clean_title, title_error = _clean_text(title, "title", HARD_TEXT_CHARS)
    if title_error:
        return None, title_error

    if not isinstance(slides, list) or not slides:
        return None, "Error: slides must be a non-empty list"
    if len(slides) > MAX_SLIDES:
        return None, f"Error: write_pptx supports at most {MAX_SLIDES} slides"

    clean_slides: list[dict[str, Any]] = []
    expected_titles: list[str] = []
    for index, raw_slide in enumerate(slides, start=1):
        clean_slide, error = _validate_slide(raw_slide, index)
        if error:
            return None, error
        assert clean_slide is not None
        clean_slides.append(clean_slide)
        slide_title = clean_slide.get("title")
        if isinstance(slide_title, str) and slide_title:
            expected_titles.append(slide_title)

    clean_theme, theme_error = _validate_theme(theme or {})
    if theme_error:
        return None, theme_error

    clean_slides, layout_warnings = _apply_layout_guardrails(clean_slides)
    expected_titles = [
        str(slide["title"]) for slide in clean_slides if isinstance(slide.get("title"), str)
    ]
    rel = _display_path(resolved, cwd)
    return (
        PresentationPlan(
            output_path=resolved,
            relative_output_path=rel,
            title=clean_title,
            slides=clean_slides,
            theme=clean_theme,
            expected_titles=expected_titles,
            overwrite=overwrite,
            layout_warnings=layout_warnings,
        ),
        None,
    )


def write_pptx(
    cwd: str,
    output_path: str,
    title: str,
    slides: list[dict[str, Any]],
    theme: dict[str, Any] | None = None,
    overwrite: bool = False,
) -> str:
    """Create a real, editable PPTX presentation and return JSON metadata."""
    plan, error = validate_pptx_request(
        cwd,
        output_path=output_path,
        title=title,
        slides=slides,
        theme=theme,
        overwrite=overwrite,
    )
    if error:
        return error
    assert plan is not None

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        return "Error: missing python-pptx dependency. Run: pip install python-pptx"

    prs = Presentation()
    prs.core_properties.title = plan.title
    prs.core_properties.subject = "Generated by CI2Lab write_pptx"
    blank_layout = prs.slide_layouts[6]
    theme_values = _theme_values(plan.theme)

    for slide_index, slide_data in enumerate(plan.slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _paint_background(slide, theme_values["background"], RGBColor)
        slide_type = slide_data["type"]
        if slide_type == "cover":
            _render_cover(slide, slide_data, theme_values, Inches, Pt, PP_ALIGN, RGBColor, MSO_SHAPE)
        elif slide_type == "section":
            _render_section(slide, slide_data, theme_values, Inches, Pt, PP_ALIGN, RGBColor, MSO_SHAPE)
        elif slide_type == "bullets":
            _render_bullets(slide, slide_data, theme_values, Inches, Pt, RGBColor)
        elif slide_type == "two_columns":
            _render_two_columns(slide, slide_data, theme_values, Inches, Pt, RGBColor)
        elif slide_type == "table":
            _render_table(slide, slide_data, theme_values, Inches, Pt, RGBColor)
        elif slide_type == "quote":
            _render_quote(slide, slide_data, theme_values, Inches, Pt, PP_ALIGN, RGBColor)
        elif slide_type == "closing":
            _render_closing(slide, slide_data, theme_values, Inches, Pt, PP_ALIGN, RGBColor)
        elif slide_type == "metric_cards":
            _render_metric_cards(slide, slide_data, theme_values, Inches, Pt, RGBColor, MSO_SHAPE)
        elif slide_type == "comparison":
            _render_comparison(slide, slide_data, theme_values, Inches, Pt, RGBColor, MSO_SHAPE)
        elif slide_type == "decision":
            _render_decision(slide, slide_data, theme_values, Inches, Pt, RGBColor, MSO_SHAPE)
        else:  # pragma: no cover - validate_pptx_request prevents this.
            return f"Error: unsupported slide type: {slide_type}"
        _add_footer_and_slide_number(
            slide,
            slide_index=slide_index,
            total_slides=len(plan.slides),
            theme=theme_values,
            Inches=Inches,
            Pt=Pt,
            RGBColor=RGBColor,
        )

    plan.output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(plan.output_path))
    except Exception as exc:
        return f"Error: could not save PPTX: {exc}"

    validation = _post_validate(plan)
    validation["layout_warnings"] = plan.layout_warnings
    if validation["status"] != "passed":
        return "Error: PPTX post-validation failed: " + "; ".join(validation["checks"])

    metadata = {
        "output_path": plan.relative_output_path,
        "written_file": plan.relative_output_path,
        "written_files": [plan.relative_output_path],
        "slide_count": len(plan.slides),
        "validation_summary": validation,
    }
    return json.dumps(metadata, ensure_ascii=False, sort_keys=True)


def _validate_slide(raw_slide: Any, index: int) -> tuple[dict[str, Any] | None, str | None]:
    if not isinstance(raw_slide, dict):
        return None, f"Error: slide {index} must be an object"
    slide_type = str(raw_slide.get("type", "")).strip().lower()
    if slide_type not in SUPPORTED_SLIDE_TYPES:
        return None, f"Error: unsupported slide type at slide {index}: {slide_type or '(missing)'}"

    clean: dict[str, Any] = {"type": slide_type}
    if slide_type in {
        "cover",
        "section",
        "bullets",
        "two_columns",
        "table",
        "closing",
        "metric_cards",
        "comparison",
        "decision",
    }:
        title, error = _required_text(raw_slide, "title", index, HARD_TEXT_CHARS)
        if error:
            return None, error
        clean["title"] = title

    if slide_type in {"cover", "section", "closing"}:
        subtitle, error = _optional_text(raw_slide, "subtitle", HARD_TEXT_CHARS)
        if error:
            return None, f"Error: slide {index} {error}"
        if subtitle:
            clean["subtitle"] = subtitle

    if slide_type in {"bullets", "closing"}:
        bullets, error = _clean_string_list(raw_slide.get("bullets", []), "bullets", index)
        if error:
            return None, error
        if slide_type == "bullets" and not bullets:
            return None, f"Error: slide {index} bullets requires a non-empty bullets list"
        clean["bullets"] = bullets

    if slide_type == "two_columns":
        left, error = _clean_string_list(raw_slide.get("left"), "left", index)
        if error:
            return None, error
        right, error = _clean_string_list(raw_slide.get("right"), "right", index)
        if error:
            return None, error
        if not left or not right:
            return None, f"Error: slide {index} two_columns requires non-empty left and right lists"
        clean["left"] = left
        clean["right"] = right
        for key in ("left_title", "right_title"):
            value, error = _optional_text(raw_slide, key, HARD_TEXT_CHARS)
            if error:
                return None, f"Error: slide {index} {error}"
            if value:
                clean[key] = value

    if slide_type == "comparison":
        left_items, error = _clean_string_list(raw_slide.get("left_items"), "left_items", index)
        if error:
            return None, error
        right_items, error = _clean_string_list(raw_slide.get("right_items"), "right_items", index)
        if error:
            return None, error
        if not left_items or not right_items:
            return None, f"Error: slide {index} comparison requires non-empty item lists"
        clean["left_items"] = left_items
        clean["right_items"] = right_items
        for key in ("left_title", "right_title"):
            value, error = _required_text(raw_slide, key, index, HARD_TEXT_CHARS)
            if error:
                return None, error
            clean[key] = value

    if slide_type == "metric_cards":
        cards, error = _clean_metric_cards(raw_slide.get("cards"), index)
        if error:
            return None, error
        clean["cards"] = cards

    if slide_type == "table":
        headers, error = _clean_string_list(raw_slide.get("headers"), "headers", index)
        if error:
            return None, error
        rows, error = _clean_table_rows(raw_slide.get("rows"), index)
        if error:
            return None, error
        if not headers:
            return None, f"Error: slide {index} table requires non-empty headers"
        if not rows:
            return None, f"Error: slide {index} table requires non-empty rows"
        if len(headers) > MAX_TABLE_COLS:
            return None, f"Error: slide {index} table supports at most {MAX_TABLE_COLS} columns"
        if any(len(row) != len(headers) for row in rows):
            return None, f"Error: slide {index} table rows must match header count"
        clean["headers"] = headers
        clean["rows"] = rows

    if slide_type == "quote":
        quote, error = _required_text(raw_slide, "quote", index, HARD_TEXT_CHARS)
        if error:
            return None, error
        clean["quote"] = quote
        for key in ("title", "author"):
            value, error = _optional_text(raw_slide, key, HARD_TEXT_CHARS)
            if error:
                return None, f"Error: slide {index} {error}"
            if value:
                clean[key] = value

    if slide_type == "decision":
        recommendation, error = _required_text(raw_slide, "recommendation", index, HARD_TEXT_CHARS)
        if error:
            return None, error
        rationale, error = _required_text(raw_slide, "rationale", index, HARD_TEXT_CHARS)
        if error:
            return None, error
        next_steps, error = _clean_string_list(raw_slide.get("next_steps"), "next_steps", index)
        if error:
            return None, error
        if not next_steps:
            return None, f"Error: slide {index} decision requires non-empty next_steps"
        clean["recommendation"] = recommendation
        clean["rationale"] = rationale
        clean["next_steps"] = next_steps

    return clean, None


def _clean_metric_cards(value: Any, index: int) -> tuple[list[dict[str, str]] | None, str | None]:
    if not isinstance(value, list) or not value:
        return None, f"Error: slide {index} metric_cards requires a non-empty cards list"
    if len(value) > MAX_CARDS:
        return None, f"Error: slide {index} metric_cards supports at most {MAX_CARDS} cards"
    cards: list[dict[str, str]] = []
    for card_index, card in enumerate(value, start=1):
        if not isinstance(card, dict):
            return None, f"Error: slide {index} card {card_index} must be an object"
        label, error = _required_text(card, "label", index, HARD_LIST_ITEM_CHARS)
        if error:
            return None, f"Error: slide {index} card {card_index} label is required"
        value_text, error = _required_text(card, "value", index, HARD_LIST_ITEM_CHARS)
        if error:
            return None, f"Error: slide {index} card {card_index} value is required"
        assert label is not None and value_text is not None
        cards.append({"label": label, "value": value_text})
    return cards, None


def _clean_table_rows(value: Any, index: int) -> tuple[list[list[str]] | None, str | None]:
    if not isinstance(value, list):
        return None, f"Error: slide {index} rows must be a list"
    if len(value) > MAX_TABLE_ROWS:
        return None, f"Error: slide {index} table supports at most {MAX_TABLE_ROWS} rows"
    rows: list[list[str]] = []
    for row_index, row in enumerate(value, start=1):
        if not isinstance(row, list):
            return None, f"Error: slide {index} row {row_index} must be a list"
        clean_row: list[str] = []
        for cell in row:
            text, error = _clean_text(str(cell), "cell", HARD_CELL_CHARS)
            if error:
                return None, f"Error: slide {index} row {row_index} {error}"
            clean_row.append(text)
        rows.append(clean_row)
    return rows, None


def _clean_string_list(value: Any, field: str, index: int) -> tuple[list[str] | None, str | None]:
    if isinstance(value, str):
        # A model sometimes emits a single string where a list of items is
        # expected (e.g. two_columns "left"). Split on newlines when present,
        # otherwise treat the whole string as one item, rather than erroring.
        value = [line for line in value.splitlines() if line.strip()] or [value]
    if not isinstance(value, list):
        return None, f"Error: slide {index} {field} must be a list"
    if len(value) > MAX_BULLETS:
        return None, f"Error: slide {index} {field} supports at most {MAX_BULLETS} items"
    clean: list[str] = []
    for item in value:
        text, error = _clean_text(str(item), field, HARD_LIST_ITEM_CHARS)
        if error:
            return None, f"Error: slide {index} {error}"
        if text:
            clean.append(text)
    return clean, None


def _required_text(
    slide: dict[str, Any],
    field: str,
    index: int,
    limit: int,
) -> tuple[str | None, str | None]:
    value, error = _optional_text(slide, field, limit)
    if error:
        return None, f"Error: slide {index} {error}"
    if not value:
        return None, f"Error: slide {index} {field} is required"
    return value, None


def _optional_text(
    slide: dict[str, Any],
    field: str,
    limit: int,
) -> tuple[str | None, str | None]:
    if field not in slide or slide.get(field) is None:
        return None, None
    return _clean_text(str(slide.get(field)), field, limit)


def _clean_text(value: str, field: str, limit: int) -> tuple[str, str | None]:
    text = " ".join(value.replace("\r", "\n").split())
    if len(text) > limit:
        return "", f"{field} exceeds {limit} characters"
    return text, None


def _apply_layout_guardrails(
    slides: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    warnings: list[dict[str, Any]] = []
    guarded: list[dict[str, Any]] = []
    for index, slide in enumerate(slides, start=1):
        current = dict(slide)
        slide_type = str(current["type"])
        if "title" in current:
            limit = COVER_TITLE_LIMIT if slide_type == "cover" else TITLE_LAYOUT_LIMIT
            current["title"] = _fit_field(
                current["title"],
                slide_index=index,
                slide_type=slide_type,
                field="title",
                limit=limit,
                warnings=warnings,
            )
        if "subtitle" in current:
            limit = COVER_SUBTITLE_LIMIT if slide_type == "cover" else BULLET_LAYOUT_LIMIT
            current["subtitle"] = _fit_field(
                current["subtitle"],
                slide_index=index,
                slide_type=slide_type,
                field="subtitle",
                limit=limit,
                warnings=warnings,
            )
        if slide_type == "cover":
            title_lines = _estimate_cover_title_lines(current.get("title", ""))
            current["_title_lines"] = title_lines
            if current.get("subtitle") and title_lines > 1:
                warnings.append(
                    _layout_warning(
                        slide_index=index,
                        slide_type=slide_type,
                        field="title",
                        length=len(current.get("title", "")),
                        limit=_cover_single_line_chars(),
                        action="cover_title_wrap_risk",
                    )
                )
        if slide_type in {"bullets", "closing"}:
            current["bullets"] = _fit_text_list(
                current.get("bullets", []),
                slide_index=index,
                slide_type=slide_type,
                field="bullets",
                item_limit=BULLET_LAYOUT_LIMIT,
                count_limit=MAX_BULLETS,
                warnings=warnings,
            )
            current["_density"] = _list_density(current["bullets"], BULLET_LAYOUT_LIMIT)
        elif slide_type == "two_columns":
            for field in ("left", "right"):
                current[field] = _fit_text_list(
                    current.get(field, []),
                    slide_index=index,
                    slide_type=slide_type,
                    field=field,
                    item_limit=COLUMN_ITEM_LIMIT,
                    count_limit=MAX_BULLETS,
                    warnings=warnings,
                )
            current["_density"] = max(
                _list_density(current["left"], COLUMN_ITEM_LIMIT),
                _list_density(current["right"], COLUMN_ITEM_LIMIT),
            )
        elif slide_type == "comparison":
            for field in ("left_items", "right_items"):
                current[field] = _fit_text_list(
                    current.get(field, []),
                    slide_index=index,
                    slide_type=slide_type,
                    field=field,
                    item_limit=COLUMN_ITEM_LIMIT,
                    count_limit=MAX_BULLETS,
                    warnings=warnings,
                )
            current["_density"] = max(
                _list_density(current["left_items"], COLUMN_ITEM_LIMIT),
                _list_density(current["right_items"], COLUMN_ITEM_LIMIT),
            )
        elif slide_type == "metric_cards":
            fitted_cards: list[dict[str, str]] = []
            for card_index, card in enumerate(current.get("cards", []), start=1):
                fitted_cards.append(
                    {
                        "label": _fit_field(
                            card["label"],
                            slide_index=index,
                            slide_type=slide_type,
                            field=f"cards[{card_index}].label",
                            limit=CARD_LABEL_LIMIT,
                            warnings=warnings,
                        ),
                        "value": _fit_field(
                            card["value"],
                            slide_index=index,
                            slide_type=slide_type,
                            field=f"cards[{card_index}].value",
                            limit=CARD_VALUE_LIMIT,
                            warnings=warnings,
                        ),
                    }
                )
            current["cards"] = fitted_cards
            current["_density"] = 1 if len(fitted_cards) > 4 else 0
        elif slide_type == "decision":
            current["recommendation"] = _fit_field(
                current["recommendation"],
                slide_index=index,
                slide_type=slide_type,
                field="recommendation",
                limit=DECISION_RECOMMENDATION_LIMIT,
                warnings=warnings,
            )
            current["rationale"] = _fit_field(
                current["rationale"],
                slide_index=index,
                slide_type=slide_type,
                field="rationale",
                limit=DECISION_RATIONALE_LIMIT,
                warnings=warnings,
            )
            current["next_steps"] = _fit_text_list(
                current["next_steps"],
                slide_index=index,
                slide_type=slide_type,
                field="next_steps",
                item_limit=COLUMN_ITEM_LIMIT,
                count_limit=MAX_BULLETS,
                warnings=warnings,
            )
            current["_density"] = _list_density(current["next_steps"], COLUMN_ITEM_LIMIT)
        elif slide_type == "table":
            current["headers"] = _fit_text_list(
                current["headers"],
                slide_index=index,
                slide_type=slide_type,
                field="headers",
                item_limit=TABLE_CELL_LAYOUT_LIMIT,
                count_limit=MAX_TABLE_COLS,
                warnings=warnings,
            )
            rows: list[list[str]] = []
            for row_index, row in enumerate(current["rows"], start=1):
                rows.append(
                    [
                        _fit_field(
                            cell,
                            slide_index=index,
                            slide_type=slide_type,
                            field=f"rows[{row_index}]",
                            limit=TABLE_CELL_LAYOUT_LIMIT,
                            warnings=warnings,
                        )
                        for cell in row
                    ]
                )
            current["rows"] = rows
            current["_density"] = 1 if len(rows) > 10 or len(current["headers"]) > 4 else 0
        elif slide_type == "quote":
            current["quote"] = _fit_field(
                current["quote"],
                slide_index=index,
                slide_type=slide_type,
                field="quote",
                limit=DECISION_RATIONALE_LIMIT,
                warnings=warnings,
            )
        guarded.append(current)
    return guarded, warnings


def _fit_text_list(
    values: list[str],
    *,
    slide_index: int,
    slide_type: str,
    field: str,
    item_limit: int,
    count_limit: int,
    warnings: list[dict[str, Any]],
) -> list[str]:
    fitted = list(values)
    if len(fitted) > count_limit:
        warnings.append(
            _layout_warning(
                slide_index=slide_index,
                slide_type=slide_type,
                field=field,
                length=len(fitted),
                limit=count_limit,
                action="truncated_items",
            )
        )
        fitted = fitted[:count_limit]
    return [
        _fit_field(
            value,
            slide_index=slide_index,
            slide_type=slide_type,
            field=f"{field}[{idx}]",
            limit=item_limit,
            warnings=warnings,
        )
        for idx, value in enumerate(fitted, start=1)
    ]


def _fit_field(
    value: str,
    *,
    slide_index: int,
    slide_type: str,
    field: str,
    limit: int,
    warnings: list[dict[str, Any]],
) -> str:
    text = " ".join(str(value).replace("\r", "\n").split())
    overflow = _estimate_overflow_risk(text, limit)
    if overflow == "high":
        warnings.append(
            _layout_warning(
                slide_index=slide_index,
                slide_type=slide_type,
                field=field,
                length=len(text),
                limit=limit,
                action="truncated_text",
            )
        )
        return _ellipsize(text, limit)
    if overflow == "medium":
        warnings.append(
            _layout_warning(
                slide_index=slide_index,
                slide_type=slide_type,
                field=field,
                length=len(text),
                limit=limit,
                action="overflow_risk",
            )
        )
    return text


def _layout_warning(
    *,
    slide_index: int,
    slide_type: str,
    field: str,
    length: int,
    limit: int,
    action: str,
) -> dict[str, Any]:
    warning = {
        "slide_index": slide_index,
        "slide_type": slide_type,
        "field": field,
        "length": length,
        "limit": limit,
        "action": action,
    }
    card_match = re.match(r"^cards\[(\d+)\]\.", field)
    if card_match:
        warning["card_index"] = int(card_match.group(1))
    return warning


def _estimate_overflow_risk(text: str, limit: int) -> str | None:
    length = len(text)
    if length > limit:
        return "high"
    if length > int(limit * 0.85):
        return "medium"
    return None


def _ellipsize(text: str, limit: int) -> str:
    if limit <= 1:
        return "…"
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


def _list_density(values: list[str], item_limit: int) -> int:
    if len(values) >= 6:
        return 2
    if len(values) >= 4 or sum(len(value) for value in values) > item_limit * 3:
        return 1
    return 0


def _validate_theme(theme: dict[str, Any]) -> tuple[dict[str, Any], str | None]:
    if not isinstance(theme, dict):
        return {}, "Error: theme must be an object when provided"
    clean: dict[str, Any] = {
        "font_family": "Aptos",
        "title_font_size": 30,
        "body_font_size": 18,
        "primary_color": "#2563EB",
        "secondary_color": "#64748B",
        "background_color": "#FFFFFF",
        "footer_text": "",
        "slide_number": False,
    }
    color_aliases = {
        "primary_color": ("primary_color", "accent_color"),
        "secondary_color": ("secondary_color", "text_color"),
        "background_color": ("background_color",),
    }
    for target, aliases in color_aliases.items():
        for alias in aliases:
            if alias not in theme or theme[alias] is None:
                continue
            value = str(theme[alias]).strip()
            if not _HEX_COLOR_RE.match(value):
                return {}, f"Error: theme.{alias} must be a hex color"
            clean[target] = value if value.startswith("#") else f"#{value}"
            break

    font_family, error = _clean_theme_string(
        theme.get("font_family", clean["font_family"]),
        "font_family",
        MAX_FONT_FAMILY_CHARS,
    )
    if error:
        return {}, error
    clean["font_family"] = font_family

    footer_text, error = _clean_theme_string(
        theme.get("footer_text", clean["footer_text"]),
        "footer_text",
        MAX_FOOTER_CHARS,
        allow_empty=True,
    )
    if error:
        return {}, error
    clean["footer_text"] = footer_text

    for key, min_size, max_size in (
        ("title_font_size", MIN_TITLE_FONT_SIZE, MAX_TITLE_FONT_SIZE),
        ("body_font_size", MIN_BODY_FONT_SIZE, MAX_BODY_FONT_SIZE),
    ):
        raw = theme.get(key, clean[key])
        if isinstance(raw, bool):
            return {}, f"Error: theme.{key} must be a number"
        try:
            size = int(raw)
        except (TypeError, ValueError):
            return {}, f"Error: theme.{key} must be a number"
        if not min_size <= size <= max_size:
            return {}, f"Error: theme.{key} must be between {min_size} and {max_size}"
        clean[key] = size

    if "slide_number" in theme:
        if not isinstance(theme["slide_number"], bool):
            return {}, "Error: theme.slide_number must be a boolean"
        clean["slide_number"] = theme["slide_number"]
    return clean, None


def _clean_theme_string(
    value: Any,
    field: str,
    limit: int,
    *,
    allow_empty: bool = False,
) -> tuple[str, str | None]:
    text, error = _clean_text(str(value or ""), field, limit)
    if error:
        return "", f"Error: theme.{error}"
    if not allow_empty and not text:
        return "", f"Error: theme.{field} is required"
    return text, None


def _theme_values(theme: dict[str, Any]) -> dict[str, Any]:
    return {
        "background": str(theme["background_color"]),
        "primary": str(theme["primary_color"]),
        "secondary": str(theme["secondary_color"]),
        "text": "#1F2937",
        "font_family": str(theme["font_family"]),
        "title_font_size": int(theme["title_font_size"]),
        "body_font_size": int(theme["body_font_size"]),
        "footer_text": str(theme["footer_text"]),
        "slide_number": bool(theme["slide_number"]),
        "table_header": "#E5E7EB",
        "panel": "#F8FAFC",
    }


def _rgb(hex_color: str, RGBColor: Any) -> Any:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _paint_background(slide: Any, color: str, RGBColor: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color, RGBColor)


def _add_rect(
    slide: Any,
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    color: str,
    RGBColor: Any,
    MSO_SHAPE: Any,
) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color, RGBColor)
    shape.line.fill.background()
    return shape


def _add_textbox(
    slide: Any,
    text: str,
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    font_size: Any,
    color: str,
    RGBColor: Any,
    font_family: str,
    bold: bool = False,
    align: Any | None = None,
) -> Any:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_family
    run.font.color.rgb = _rgb(color, RGBColor)
    return box


def _add_bullet_list(
    slide: Any,
    bullets: list[str],
    *,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    font_size: Any,
    color: str,
    RGBColor: Any,
    font_family: str,
) -> Any:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = font_size
        paragraph.font.name = font_family
        paragraph.font.color.rgb = _rgb(color, RGBColor)
    return box


def _body_font_size(theme: dict[str, Any], *, density: int = 0, delta: int = 0) -> int:
    return max(MIN_FITTED_BODY_FONT_SIZE, int(theme["body_font_size"]) + delta - density * 2)


def _title_font_size(theme: dict[str, Any], *, delta: int = 0) -> int:
    return max(MIN_TITLE_FONT_SIZE, int(theme["title_font_size"]) + delta)


def _cover_single_line_chars() -> int:
    return 28


def _estimate_cover_title_lines(title: str) -> int:
    if not title:
        return 1
    longest_word = max((len(word) for word in title.split()), default=0)
    weighted_length = len(title) + max(0, longest_word - 14)
    return max(1, min(3, (weighted_length + _cover_single_line_chars() - 1) // _cover_single_line_chars()))


def _cover_title_font_size(title: str, theme: dict[str, Any]) -> int:
    lines = _estimate_cover_title_lines(title)
    base_size = max(int(theme["title_font_size"]) + 10, 40)
    if lines > 1:
        base_size -= (lines - 1) * 6
    if len(title) > 70:
        base_size -= 4
    return max(MIN_COVER_TITLE_FONT_SIZE, base_size)


def _cover_subtitle_font_size(subtitle: str, theme: dict[str, Any], *, title_lines: int) -> int:
    size = int(theme["body_font_size"]) + 2
    if len(subtitle) > 85:
        size -= 2
    if title_lines >= 3:
        size -= 2
    return max(MIN_COVER_SUBTITLE_FONT_SIZE, size)


def _metric_value_font_size(value: str, theme: dict[str, Any], *, density: int = 0) -> int:
    size = int(theme["title_font_size"]) - 8 - density * 2
    if len(value) > 18:
        size -= 3
    if len(value) > 24:
        size -= 2
    if len(value.split()) >= 3:
        size -= 2
    return max(14, size)


def _metric_label_font_size(theme: dict[str, Any], *, density: int = 0) -> int:
    return max(9, min(12, _body_font_size(theme, density=density, delta=-6)))


def _render_cover(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, PP_ALIGN: Any, RGBColor: Any, MSO_SHAPE: Any) -> None:
    _add_rect(slide, left=Inches(0.9), top=Inches(1.75), width=Inches(8.2), height=Inches(0.08), color=theme["primary"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
    title_lines = int(data.get("_title_lines") or _estimate_cover_title_lines(data["title"]))
    title_top = 1.92
    title_height = min(2.2, 0.74 * title_lines + 0.55)
    subtitle_gap = 0.22
    _add_textbox(slide, data["title"], left=Inches(0.8), top=Inches(title_top), width=Inches(8.4), height=Inches(title_height), font_size=Pt(_cover_title_font_size(data["title"], theme)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True, align=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        subtitle_top = title_top + title_height + subtitle_gap
        subtitle_height = 0.95 if len(data["subtitle"]) > 85 else 0.75
        _add_textbox(slide, data["subtitle"], left=Inches(1.1), top=Inches(subtitle_top), width=Inches(7.8), height=Inches(subtitle_height), font_size=Pt(_cover_subtitle_font_size(data["subtitle"], theme, title_lines=title_lines)), color=theme["secondary"], RGBColor=RGBColor, font_family=theme["font_family"], align=PP_ALIGN.CENTER)


def _render_section(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, PP_ALIGN: Any, RGBColor: Any, MSO_SHAPE: Any) -> None:
    _add_rect(slide, left=Inches(0), top=Inches(0), width=Inches(0.22), height=Inches(7.5), color=theme["primary"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
    _add_rect(slide, left=Inches(0.22), top=Inches(0), width=Inches(9.78), height=Inches(7.5), color=theme["panel"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
    _add_textbox(slide, data["title"], left=Inches(0.8), top=Inches(2.2), width=Inches(8.4), height=Inches(0.9), font_size=Pt(theme["title_font_size"] + 4), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True, align=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        _add_textbox(slide, data["subtitle"], left=Inches(1.2), top=Inches(3.2), width=Inches(7.6), height=Inches(0.7), font_size=Pt(theme["body_font_size"]), color=theme["secondary"], RGBColor=RGBColor, font_family=theme["font_family"], align=PP_ALIGN.CENTER)


def _render_bullets(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any) -> None:
    density = int(data.get("_density", 0))
    _add_textbox(slide, data["title"], left=Inches(0.55), top=Inches(0.4), width=Inches(9.0), height=Inches(0.7), font_size=Pt(_title_font_size(theme)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_bullet_list(slide, data["bullets"], left=Inches(0.85), top=Inches(1.45), width=Inches(8.5), height=Inches(4.6), font_size=Pt(_body_font_size(theme, density=density)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])


def _render_two_columns(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any) -> None:
    density = int(data.get("_density", 0))
    _add_textbox(slide, data["title"], left=Inches(0.55), top=Inches(0.35), width=Inches(9.0), height=Inches(0.6), font_size=Pt(_title_font_size(theme, delta=-2)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    if data.get("left_title"):
        _add_textbox(slide, data["left_title"], left=Inches(0.75), top=Inches(1.25), width=Inches(4.0), height=Inches(0.4), font_size=Pt(_body_font_size(theme, density=density)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    if data.get("right_title"):
        _add_textbox(slide, data["right_title"], left=Inches(5.15), top=Inches(1.25), width=Inches(4.0), height=Inches(0.4), font_size=Pt(_body_font_size(theme, density=density)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_bullet_list(slide, data["left"], left=Inches(0.75), top=Inches(1.8), width=Inches(4.0), height=Inches(4.3), font_size=Pt(_body_font_size(theme, density=density, delta=-1)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])
    _add_bullet_list(slide, data["right"], left=Inches(5.15), top=Inches(1.8), width=Inches(4.0), height=Inches(4.3), font_size=Pt(_body_font_size(theme, density=density, delta=-1)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])


def _render_table(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any) -> None:
    _add_textbox(slide, data["title"], left=Inches(0.55), top=Inches(0.35), width=Inches(9.0), height=Inches(0.6), font_size=Pt(theme["title_font_size"] - 2), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    row_count = len(data["rows"]) + 1
    col_count = len(data["headers"])
    table_shape = slide.shapes.add_table(row_count, col_count, Inches(0.65), Inches(1.35), Inches(8.7), Inches(4.6))
    table = table_shape.table
    for col_idx, header in enumerate(data["headers"]):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(theme["table_header"], RGBColor)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.name = theme["font_family"]
    for row_idx, row in enumerate(data["rows"], start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = theme["font_family"]


def _render_quote(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, PP_ALIGN: Any, RGBColor: Any) -> None:
    if data.get("title"):
        _add_textbox(slide, data["title"], left=Inches(0.55), top=Inches(0.4), width=Inches(9.0), height=Inches(0.6), font_size=Pt(theme["title_font_size"] - 4), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_textbox(slide, data["quote"], left=Inches(1.0), top=Inches(2.0), width=Inches(8.0), height=Inches(1.4), font_size=Pt(theme["title_font_size"] - 2), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"], align=PP_ALIGN.CENTER)
    if data.get("author"):
        _add_textbox(slide, data["author"], left=Inches(1.0), top=Inches(3.6), width=Inches(8.0), height=Inches(0.5), font_size=Pt(max(theme["body_font_size"] - 2, MIN_BODY_FONT_SIZE)), color=theme["secondary"], RGBColor=RGBColor, font_family=theme["font_family"], align=PP_ALIGN.CENTER)


def _render_closing(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, PP_ALIGN: Any, RGBColor: Any) -> None:
    _add_textbox(slide, data["title"], left=Inches(0.8), top=Inches(1.6), width=Inches(8.4), height=Inches(0.9), font_size=Pt(theme["title_font_size"] + 4), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True, align=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        _add_textbox(slide, data["subtitle"], left=Inches(1.0), top=Inches(2.6), width=Inches(8.0), height=Inches(0.6), font_size=Pt(theme["body_font_size"]), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"], align=PP_ALIGN.CENTER)
    if data.get("bullets"):
        _add_bullet_list(slide, data["bullets"], left=Inches(2.0), top=Inches(3.45), width=Inches(6.0), height=Inches(1.8), font_size=Pt(max(theme["body_font_size"] - 2, MIN_BODY_FONT_SIZE)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])


def _render_metric_cards(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any, MSO_SHAPE: Any) -> None:
    density = int(data.get("_density", 0))
    _add_textbox(slide, data["title"], left=Inches(0.55), top=Inches(0.35), width=Inches(9.0), height=Inches(0.6), font_size=Pt(_title_font_size(theme, delta=-2)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    cards = data["cards"]
    col_count = 3 if len(cards) > 3 else len(cards)
    card_width = 2.65 if col_count == 3 else 8.4 / max(col_count, 1)
    card_height = 1.45
    row_gap = 1.85
    for idx, card in enumerate(cards):
        row = idx // 3
        col = idx % 3
        left = Inches(0.75 + col * 3.0)
        top = Inches(1.25 + row * row_gap)
        _add_rect(slide, left=left, top=top, width=Inches(card_width), height=Inches(card_height), color=theme["panel"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
        _add_rect(slide, left=left, top=top, width=Inches(0.08), height=Inches(card_height), color=theme["primary"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
        text_left = left + Inches(0.18)
        text_width = Inches(card_width - 0.32)
        label_top = top + Inches(0.17)
        value_top = top + Inches(0.55)
        _add_textbox(slide, card["label"], left=text_left, top=label_top, width=text_width, height=Inches(0.28), font_size=Pt(_metric_label_font_size(theme, density=density)), color=theme["secondary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
        _add_textbox(slide, card["value"], left=text_left, top=value_top, width=text_width, height=Inches(0.72), font_size=Pt(_metric_value_font_size(card["value"], theme, density=density)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)


def _render_comparison(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any, MSO_SHAPE: Any) -> None:
    density = int(data.get("_density", 0))
    _add_textbox(slide, data["title"], left=Inches(0.55), top=Inches(0.35), width=Inches(9.0), height=Inches(0.6), font_size=Pt(_title_font_size(theme, delta=-2)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_rect(slide, left=Inches(0.7), top=Inches(1.25), width=Inches(4.1), height=Inches(4.85), color=theme["panel"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
    _add_rect(slide, left=Inches(5.1), top=Inches(1.25), width=Inches(4.1), height=Inches(4.85), color=theme["panel"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
    _add_textbox(slide, data["left_title"], left=Inches(0.95), top=Inches(1.5), width=Inches(3.6), height=Inches(0.4), font_size=Pt(_body_font_size(theme, density=density, delta=1)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_textbox(slide, data["right_title"], left=Inches(5.35), top=Inches(1.5), width=Inches(3.6), height=Inches(0.4), font_size=Pt(_body_font_size(theme, density=density, delta=1)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_bullet_list(slide, data["left_items"], left=Inches(0.95), top=Inches(2.05), width=Inches(3.55), height=Inches(3.5), font_size=Pt(_body_font_size(theme, density=density, delta=-2)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])
    _add_bullet_list(slide, data["right_items"], left=Inches(5.35), top=Inches(2.05), width=Inches(3.55), height=Inches(3.5), font_size=Pt(_body_font_size(theme, density=density, delta=-2)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])


def _render_decision(slide: Any, data: dict[str, Any], theme: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any, MSO_SHAPE: Any) -> None:
    density = int(data.get("_density", 0))
    _add_textbox(slide, data["title"], left=Inches(0.55), top=Inches(0.35), width=Inches(9.0), height=Inches(0.6), font_size=Pt(_title_font_size(theme, delta=-2)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_rect(slide, left=Inches(0.75), top=Inches(1.25), width=Inches(8.5), height=Inches(1.15), color=theme["primary"], RGBColor=RGBColor, MSO_SHAPE=MSO_SHAPE)
    _add_textbox(slide, data["recommendation"], left=Inches(0.95), top=Inches(1.48), width=Inches(8.05), height=Inches(0.55), font_size=Pt(_body_font_size(theme, density=density, delta=2)), color="#FFFFFF", RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_textbox(slide, data["rationale"], left=Inches(0.85), top=Inches(2.75), width=Inches(8.3), height=Inches(1.1), font_size=Pt(_body_font_size(theme, density=density)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])
    _add_textbox(slide, "Siguientes pasos", left=Inches(0.85), top=Inches(4.05), width=Inches(8.3), height=Inches(0.35), font_size=Pt(_body_font_size(theme, density=density)), color=theme["primary"], RGBColor=RGBColor, font_family=theme["font_family"], bold=True)
    _add_bullet_list(slide, data["next_steps"], left=Inches(1.05), top=Inches(4.5), width=Inches(7.8), height=Inches(1.45), font_size=Pt(_body_font_size(theme, density=density, delta=-2)), color=theme["text"], RGBColor=RGBColor, font_family=theme["font_family"])


def _add_footer_and_slide_number(
    slide: Any,
    *,
    slide_index: int,
    total_slides: int,
    theme: dict[str, Any],
    Inches: Any,
    Pt: Any,
    RGBColor: Any,
) -> None:
    footer = theme.get("footer_text", "")
    if footer:
        _add_textbox(slide, footer, left=Inches(0.55), top=Inches(6.95), width=Inches(6.8), height=Inches(0.25), font_size=Pt(8), color=theme["secondary"], RGBColor=RGBColor, font_family=theme["font_family"])
    if theme.get("slide_number"):
        _add_textbox(slide, f"{slide_index}/{total_slides}", left=Inches(8.75), top=Inches(6.95), width=Inches(0.7), height=Inches(0.25), font_size=Pt(8), color=theme["secondary"], RGBColor=RGBColor, font_family=theme["font_family"])


def _post_validate(plan: PresentationPlan) -> dict[str, Any]:
    checks: list[str] = []
    try:
        from pptx import Presentation
    except ImportError:
        return {"status": "failed", "checks": ["python-pptx import failed during validation"]}

    if not plan.output_path.is_file():
        checks.append("output file exists: failed")
    elif plan.output_path.stat().st_size <= 0:
        checks.append("output file non-empty: failed")
    else:
        checks.append("output file non-empty: passed")

    try:
        reopened = Presentation(str(plan.output_path))
    except Exception as exc:
        return {"status": "failed", "checks": [*checks, f"reopen pptx: failed ({exc})"]}

    if len(reopened.slides) == len(plan.slides):
        checks.append("slide count matches: passed")
    else:
        checks.append(f"slide count matches: failed ({len(reopened.slides)} != {len(plan.slides)})")

    all_text = "\n".join(_shape_text(shape) for slide in reopened.slides for shape in slide.shapes)
    missing = [title for title in plan.expected_titles if title not in all_text]
    if missing:
        checks.append("expected titles present: failed (" + ", ".join(missing) + ")")
    else:
        checks.append("expected titles present: passed")

    status = "failed" if any("failed" in check for check in checks) else "passed"
    return {"status": status, "checks": checks}


def _shape_text(shape: Any) -> str:
    if hasattr(shape, "text"):
        return str(shape.text)
    if hasattr(shape, "table"):
        values: list[str] = []
        for row in shape.table.rows:
            for cell in row.cells:
                values.append(str(cell.text))
        return "\n".join(values)
    return ""


def _display_path(resolved: Path, cwd: str) -> str:
    try:
        return str(resolved.relative_to(Path(cwd).resolve()))
    except ValueError:
        return str(resolved)
