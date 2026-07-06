"""Tests for write_pptx."""

from __future__ import annotations

import json
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

from ci2lab.harness.tools.pptx_writer import write_pptx
from ci2lab.harness.tools.registry import execute_tool
from ci2lab.harness.types import AgentConfig, ToolCall


def _metadata(result: str) -> dict:
    assert not result.startswith("Error:"), result
    return json.loads(result)


def _all_text(path: Path) -> str:
    prs = Presentation(str(path))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(str(shape.text))
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(str(cell.text))
    return "\n".join(chunks)


def _first_run_font(path: Path) -> str | None:
    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        return run.font.name
    return None


def _text_shape_bounds(path: Path, text: str) -> tuple[int, int, int, int]:
    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text == text:
                return shape.left, shape.top, shape.width, shape.height
    raise AssertionError(f"Text shape not found: {text}")


def _assert_textbox_above(path: Path, upper_text: str, lower_text: str) -> None:
    _, upper_top, _, upper_height = _text_shape_bounds(path, upper_text)
    _, lower_top, _, _ = _text_shape_bounds(path, lower_text)
    assert upper_top + upper_height <= lower_top


def test_write_pptx_minimal_cover_and_bullets(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "deck.pptx",
        "Demo deck",
        [
            {"type": "cover", "title": "Demo deck", "subtitle": "CI2Lab"},
            {"type": "bullets", "title": "Highlights", "bullets": ["One", "Two"]},
        ],
    )

    meta = _metadata(result)
    assert meta["output_path"] == "deck.pptx"
    assert meta["written_file"] == "deck.pptx"
    assert meta["written_files"] == ["deck.pptx"]
    assert meta["slide_count"] == 2
    assert meta["validation_summary"]["status"] == "passed"
    assert len(Presentation(str(tmp_path / "deck.pptx")).slides) == 2
    text = _all_text(tmp_path / "deck.pptx")
    assert "Demo deck" in text
    assert "Highlights" in text


def test_write_pptx_cover_short_title_has_no_layout_warnings(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "cover-short.pptx",
        "Cover",
        [{"type": "cover", "title": "Demo deck", "subtitle": "CI2Lab"}],
    )

    meta = _metadata(result)
    assert meta["validation_summary"]["layout_warnings"] == []
    assert len(Presentation(str(tmp_path / "cover-short.pptx")).slides) == 1
    _assert_textbox_above(tmp_path / "cover-short.pptx", "Demo deck", "CI2Lab")


def test_write_pptx_cover_long_title_keeps_subtitle_below_title(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "cover-long.pptx",
        "Cover",
        [
            {
                "type": "cover",
                "title": "Hardware local para LLMs grandes",
                "subtitle": "RUDIGER, modelos 32B/70B y el umbral practico de 405B",
            }
        ],
        theme={"title_font_size": 34, "body_font_size": 18},
    )

    meta = _metadata(result)
    warnings = meta["validation_summary"]["layout_warnings"]
    assert any(warning["action"] == "cover_title_wrap_risk" for warning in warnings)
    assert len(Presentation(str(tmp_path / "cover-long.pptx")).slides) == 1
    _assert_textbox_above(
        tmp_path / "cover-long.pptx",
        "Hardware local para LLMs grandes",
        "RUDIGER, modelos 32B/70B y el umbral practico de 405B",
    )


def test_write_pptx_cover_very_long_title_generates_layout_warning(tmp_path: Path) -> None:
    long_title = (
        "Hardware local para inferencia de LLMs grandes con presupuesto, "
        "latencia, memoria y comparacion cloud"
    )
    result = write_pptx(
        str(tmp_path),
        "cover-very-long.pptx",
        "Cover",
        [{"type": "cover", "title": long_title, "subtitle": "Decision tecnica"}],
    )

    meta = _metadata(result)
    warnings = meta["validation_summary"]["layout_warnings"]
    assert any(warning["field"] == "title" for warning in warnings)
    assert len(Presentation(str(tmp_path / "cover-very-long.pptx")).slides) == 1
    assert "Decision tecnica" in _all_text(tmp_path / "cover-very-long.pptx")


def test_write_pptx_two_columns(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "columns.pptx",
        "Columns",
        [
            {
                "type": "two_columns",
                "title": "Compare",
                "left_title": "Pros",
                "right_title": "Cons",
                "left": ["Fast", "Simple"],
                "right": ["Limited"],
            }
        ],
    )

    _metadata(result)
    text = _all_text(tmp_path / "columns.pptx")
    assert "Compare" in text
    assert "Pros" in text
    assert "Cons" in text
    assert "Fast" in text


def test_write_pptx_table(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "table.pptx",
        "Table",
        [
            {
                "type": "table",
                "title": "Results",
                "headers": ["Metric", "Value"],
                "rows": [["Accuracy", "95%"], ["Latency", "Low"]],
            }
        ],
    )

    _metadata(result)
    text = _all_text(tmp_path / "table.pptx")
    assert "Results" in text
    assert "Metric" in text
    assert "Accuracy" in text


def test_write_pptx_rejects_path_outside_workspace(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "../outside.pptx",
        "Bad",
        [{"type": "cover", "title": "Bad"}],
    )

    assert result.startswith("Error:")
    assert "outside" in result.lower() or "escapes" in result.lower()


def test_write_pptx_rejects_invalid_extension(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "deck.txt",
        "Bad",
        [{"type": "cover", "title": "Bad"}],
    )

    assert result == "Error: write_pptx only accepts .pptx output paths"


def test_write_pptx_rejects_unknown_slide_type(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "deck.pptx",
        "Bad",
        [{"type": "agenda", "title": "Agenda"}],
    )

    assert result.startswith("Error: unsupported slide type")


def test_write_pptx_rejects_empty_slides(tmp_path: Path) -> None:
    result = write_pptx(str(tmp_path), "deck.pptx", "Bad", [])

    assert result == "Error: slides must be a non-empty list"


def test_write_pptx_rejects_overwrite_without_flag(tmp_path: Path) -> None:
    target = tmp_path / "deck.pptx"
    first = write_pptx(str(tmp_path), "deck.pptx", "Deck", [{"type": "cover", "title": "Deck"}])
    assert not first.startswith("Error:")

    result = write_pptx(
        str(tmp_path),
        "deck.pptx",
        "Deck 2",
        [{"type": "cover", "title": "Deck 2"}],
    )

    assert target.is_file()
    assert result == "Error: output file already exists; set overwrite=true to replace it"


def test_execute_write_pptx_with_approval_creates_openable_file(tmp_path: Path) -> None:
    cfg = AgentConfig(
        cwd=str(tmp_path),
        security_engine="ci2lab",
        require_diff_preview=True,
        confirm_callback=lambda _name, _summary: True,
    )
    call = ToolCall(
        name="write_pptx",
        arguments={
            "output_path": "via-tool.pptx",
            "title": "Via tool",
            "slides": [{"type": "cover", "title": "Via tool"}],
        },
    )

    with patch("ci2lab.console.console.print"):
        result = execute_tool(call, cfg)

    assert not result.is_error
    assert result.outcome == "approved"
    assert len(Presentation(str(tmp_path / "via-tool.pptx")).slides) == 1


def test_write_pptx_expected_titles_appear(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "titles.pptx",
        "Titles",
        [
            {"type": "section", "title": "Context"},
            {"type": "quote", "title": "Principle", "quote": "Keep it simple."},
            {"type": "closing", "title": "Next steps"},
        ],
    )

    meta = _metadata(result)
    assert meta["validation_summary"]["status"] == "passed"
    text = _all_text(tmp_path / "titles.pptx")
    assert "Context" in text
    assert "Principle" in text
    assert "Next steps" in text


def test_write_pptx_with_complete_theme(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "theme.pptx",
        "Theme",
        [{"type": "cover", "title": "Theme", "subtitle": "Styled"}],
        theme={
            "font_family": "Arial",
            "title_font_size": 34,
            "body_font_size": 19,
            "primary_color": "#123456",
            "secondary_color": "#654321",
            "background_color": "#F7F8FA",
            "footer_text": "CI2Lab",
            "slide_number": True,
        },
    )

    meta = _metadata(result)
    assert meta["written_files"] == ["theme.pptx"]
    text = _all_text(tmp_path / "theme.pptx")
    assert "Theme" in text
    assert "CI2Lab" in text
    assert "1/1" in text
    assert _first_run_font(tmp_path / "theme.pptx") == "Arial"


def test_write_pptx_without_theme_uses_defaults(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "defaults.pptx",
        "Defaults",
        [{"type": "cover", "title": "Defaults"}],
    )

    _metadata(result)
    assert len(Presentation(str(tmp_path / "defaults.pptx")).slides) == 1
    assert _first_run_font(tmp_path / "defaults.pptx") == "Aptos"


def test_write_pptx_rejects_invalid_theme_color(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "bad-color.pptx",
        "Bad",
        [{"type": "cover", "title": "Bad"}],
        theme={"primary_color": "blue"},
    )

    assert result == "Error: theme.primary_color must be a hex color"


def test_write_pptx_rejects_theme_font_size_out_of_range(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "bad-font.pptx",
        "Bad",
        [{"type": "cover", "title": "Bad"}],
        theme={"title_font_size": 100},
    )

    assert result == "Error: theme.title_font_size must be between 18 and 60"


def test_write_pptx_metric_cards(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "metrics.pptx",
        "Metrics",
        [
            {
                "type": "metric_cards",
                "title": "Hardware snapshot",
                "cards": [
                    {"label": "RAM", "value": "128 GB"},
                    {"label": "VRAM", "value": "96 GB"},
                    {"label": "GPUs", "value": "2x A6000"},
                ],
            }
        ],
    )

    meta = _metadata(result)
    assert meta["validation_summary"]["layout_warnings"] == []
    text = _all_text(tmp_path / "metrics.pptx")
    assert "Hardware snapshot" in text
    assert "128 GB" in text
    assert "VRAM" in text
    _assert_textbox_above(tmp_path / "metrics.pptx", "RAM", "128 GB")


def test_write_pptx_metric_cards_long_value_keeps_label_above_value(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "threadripper-card.pptx",
        "Metrics",
        [
            {
                "type": "metric_cards",
                "title": "Hardware snapshot",
                "cards": [{"label": "CPU", "value": "Threadripper PRO 7975WX"}],
            }
        ],
        theme={"title_font_size": 34, "body_font_size": 18},
    )

    meta = _metadata(result)
    warnings = meta["validation_summary"]["layout_warnings"]
    assert any(warning["field"] == "cards[1].value" for warning in warnings)
    assert any(warning.get("card_index") == 1 for warning in warnings)
    assert len(Presentation(str(tmp_path / "threadripper-card.pptx")).slides) == 1
    assert "Threadripper PRO 7975WX" in _all_text(tmp_path / "threadripper-card.pptx")
    _assert_textbox_above(tmp_path / "threadripper-card.pptx", "CPU", "Threadripper PRO 7975WX")


def test_write_pptx_metric_cards_multiword_value_keeps_label_above_value(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "multiword-card.pptx",
        "Metrics",
        [
            {
                "type": "metric_cards",
                "title": "Cost scenarios",
                "cards": [{"label": "4xH100", "value": "experimental / posible"}],
            }
        ],
    )

    _metadata(result)
    assert "experimental / posible" in _all_text(tmp_path / "multiword-card.pptx")
    _assert_textbox_above(tmp_path / "multiword-card.pptx", "4xH100", "experimental / posible")


def test_write_pptx_metric_cards_six_cards_stays_openable_without_overlap(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "six-cards.pptx",
        "Metrics",
        [
            {
                "type": "metric_cards",
                "title": "Hardware snapshot",
                "cards": [
                    {"label": "CPU", "value": "Threadripper PRO 7975WX"},
                    {"label": "RAM", "value": "128 GB"},
                    {"label": "GPU", "value": "2x RTX A6000"},
                    {"label": "Sistema", "value": "Ubuntu 22.04"},
                    {"label": "4xH100", "value": "experimental / posible"},
                    {"label": "8xH200", "value": "datacenter potente"},
                ],
            }
        ],
    )

    _metadata(result)
    assert len(Presentation(str(tmp_path / "six-cards.pptx")).slides) == 1
    text = _all_text(tmp_path / "six-cards.pptx")
    assert "Hardware snapshot" in text
    assert "datacenter potente" in text
    _assert_textbox_above(tmp_path / "six-cards.pptx", "CPU", "Threadripper PRO 7975WX")
    _assert_textbox_above(tmp_path / "six-cards.pptx", "4xH100", "experimental / posible")


def test_write_pptx_comparison(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "comparison.pptx",
        "Comparison",
        [
            {
                "type": "comparison",
                "title": "Cloud vs Local",
                "left_title": "Cloud",
                "left_items": ["Fast start", "Pay per use"],
                "right_title": "Local",
                "right_items": ["Control", "Reusable"],
            }
        ],
    )

    _metadata(result)
    text = _all_text(tmp_path / "comparison.pptx")
    assert "Cloud vs Local" in text
    assert "Fast start" in text
    assert "Reusable" in text


def test_write_pptx_decision(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "decision.pptx",
        "Decision",
        [
            {
                "type": "decision",
                "title": "Recommendation",
                "recommendation": "Benchmark before buying hardware.",
                "rationale": "A short cloud campaign reduces CAPEX risk.",
                "next_steps": ["Run CI2Lab tasks", "Compare cost", "Decide"],
            }
        ],
    )

    _metadata(result)
    text = _all_text(tmp_path / "decision.pptx")
    assert "Recommendation" in text
    assert "Benchmark before buying hardware." in text
    assert "Siguientes pasos" in text


def test_write_pptx_long_bullets_generate_layout_warnings(tmp_path: Path) -> None:
    long_bullet = (
        "Este bullet conserva acentos y unicode, pero es deliberadamente largo para "
        "forzar compactación determinista dentro de la caja de texto sin romper el PPTX."
    )
    result = write_pptx(
        str(tmp_path),
        "long-bullets.pptx",
        "Long bullets",
        [{"type": "bullets", "title": "Contenido denso", "bullets": [long_bullet] * 5}],
    )

    meta = _metadata(result)
    warnings = meta["validation_summary"]["layout_warnings"]
    assert warnings
    assert any(warning["field"].startswith("bullets") for warning in warnings)
    assert len(Presentation(str(tmp_path / "long-bullets.pptx")).slides) == 1
    assert "…" in _all_text(tmp_path / "long-bullets.pptx")


def test_write_pptx_long_metric_cards_generate_layout_warnings(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "long-cards.pptx",
        "Cards",
        [
            {
                "type": "metric_cards",
                "title": "Métricas",
                "cards": [
                    {
                        "label": "Memoria de GPU disponible para experimentos locales",
                        "value": "Noventa y seis gigabytes de VRAM total útil",
                    }
                ],
            }
        ],
    )

    meta = _metadata(result)
    warnings = meta["validation_summary"]["layout_warnings"]
    assert any("cards[1]" in warning["field"] for warning in warnings)
    assert len(Presentation(str(tmp_path / "long-cards.pptx")).slides) == 1


def test_write_pptx_long_comparison_generates_layout_warnings(tmp_path: Path) -> None:
    item = (
        "Una explicación larga con acentos sobre coste, latencia, throughput, "
        "operación, mantenimiento, riesgo de infrautilización y toma de decisiones"
    )
    result = write_pptx(
        str(tmp_path),
        "long-comparison.pptx",
        "Comparison",
        [
            {
                "type": "comparison",
                "title": "Comparación",
                "left_title": "Local",
                "left_items": [item] * 4,
                "right_title": "Cloud",
                "right_items": [item] * 4,
            }
        ],
    )

    meta = _metadata(result)
    warnings = meta["validation_summary"]["layout_warnings"]
    assert any("left_items" in warning["field"] for warning in warnings)
    assert len(Presentation(str(tmp_path / "long-comparison.pptx")).slides) == 1


def test_write_pptx_long_decision_generates_layout_warnings(tmp_path: Path) -> None:
    result = write_pptx(
        str(tmp_path),
        "long-decision.pptx",
        "Decision",
        [
            {
                "type": "decision",
                "title": "Decisión",
                "recommendation": "Comprar hardware solo si la campaña cloud demuestra mejoras repetibles y económicamente defendibles con las mismas tareas de CI2Lab.",
                "rationale": (
                    "La razón principal es que el coste de capital, mantenimiento, "
                    "electricidad y administración puede superar con facilidad el valor "
                    "incremental del modelo si la tasa de éxito no mejora lo suficiente."
                ),
                "next_steps": [
                    "Ejecutar benchmark con 32B.",
                    "Ejecutar benchmark con 70B reviewer.",
                    "Comparar contra 405B cloud.",
                    "Calcular coste por tarea resuelta.",
                    "Pedir ofertas solo si los datos lo justifican.",
                ],
            }
        ],
    )

    meta = _metadata(result)
    warnings = meta["validation_summary"]["layout_warnings"]
    assert any(warning["field"] in {"recommendation", "rationale"} for warning in warnings)
    assert len(Presentation(str(tmp_path / "long-decision.pptx")).slides) == 1


def test_write_pptx_rejects_table_with_too_many_rows(tmp_path: Path) -> None:
    rows = [[str(i), "value"] for i in range(25)]
    result = write_pptx(
        str(tmp_path),
        "too-large-table.pptx",
        "Table",
        [{"type": "table", "title": "Too large", "headers": ["A", "B"], "rows": rows}],
    )

    assert result == "Error: slide 1 table supports at most 20 rows"
